#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
自工程检查重点检查线体工具 v2.0
- 读取生产计划PPT
- 根据过去40天数据计算各线体主力产品
- 标记产品变化线体为重点检查对象
- 支持导出PDF、发送邮件
"""

import os
import sys
import re
import smtplib
import json
import subprocess
from datetime import datetime, timedelta
from collections import defaultdict, Counter
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLabel, QTableWidget, QTableWidgetItem, QFileDialog,
    QMessageBox, QHeaderView, QAbstractItemView, QStatusBar, QProgressDialog,
    QDialog, QFormLayout, QLineEdit, QTextEdit, QDialogButtonBox
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QColor, QFont

from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# ==================== 邮件配置 ====================

EMAIL_CONFIG_FILE = os.path.expanduser("~/.ipqa_email_config.json")


def load_email_config():
    """加载邮件配置"""
    if os.path.exists(EMAIL_CONFIG_FILE):
        with open(EMAIL_CONFIG_FILE, 'r') as f:
            return json.load(f)
    return {
        "smtp_server": "smtp.qq.com",
        "smtp_port": 465,
        "sender": "16835577@qq.com",
        "password": "",      # 留空，用户需要填写授权码
        "recipient": ""
    }


def save_email_config(config):
    """保存邮件配置"""
    with open(EMAIL_CONFIG_FILE, 'w') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


# ==================== PPT解析逻辑 ====================

EXCLUDE_KEYWORDS = ['社内', '亏产', '艾航', '中兴', '永信', '余裕人员']
EXCLUDE_LINE_KEYWORDS = ['键盘', '踏板']


def is_valid_product(product: str) -> bool:
    if not product or not product.strip():
        return False
    product = product.strip()
    for kw in EXCLUDE_KEYWORDS:
        if kw in product:
            return False
    if re.match(r'^[\d\-\/]+$', product):
        return False
    return True


def is_valid_line(line: str) -> bool:
    if not line or not line.strip():
        return False
    for kw in EXCLUDE_LINE_KEYWORDS:
        if kw in line:
            return False
    return True


def parse_ppt(filepath: str) -> dict:
    result = {}
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = list(table.rows)
                    if len(rows) < 2:
                        continue
                    header = [cell.text.strip() for cell in rows[0].cells]
                    try:
                        line_idx = header.index('线体')
                        product_idx = header.index('当日生产机种')
                    except ValueError:
                        continue
                    for row in rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if len(cells) <= max(line_idx, product_idx):
                            continue
                        line = cells[line_idx]
                        product = cells[product_idx]
                        if is_valid_line(line) and is_valid_product(product):
                            products = [p.strip() for p in product.split('/') if is_valid_product(p.strip())]
                            if products:
                                result[line] = products
    except Exception as e:
        print(f"解析PPT出错 {filepath}: {e}")
    return result


def extract_date_from_filename(filename: str) -> datetime:
    match = re.search(r'(\d+)月(\d+)日', filename)
    if match:
        month = int(match.group(1))
        day = int(match.group(2))
        return datetime(2026, month, day)
    return datetime(2026, 1, 1)


def get_ppt_files(folder: str, days: int = 40) -> list:
    files = []
    if not os.path.exists(folder):
        return files
    today = datetime.now()
    cutoff = today - timedelta(days=days)
    for f in os.listdir(folder):
        if f.endswith('.pptx'):
            filepath = os.path.join(folder, f)
            file_date = extract_date_from_filename(f)
            if file_date >= cutoff:
                files.append((file_date, filepath, f))
    files.sort(key=lambda x: x[0])
    return files


def calculate_main_products(ppt_files: list) -> dict:
    """计算每条线体的前3名主力机型"""
    line_product_count = defaultdict(Counter)
    for _, filepath, _ in ppt_files:
        daily_data = parse_ppt(filepath)
        for line, products in daily_data.items():
            for product in products:
                line_product_count[line][product] += 1
    main_products = {}
    for line, counter in line_product_count.items():
        if counter:
            # 取频次最高的前3个产品作为主力机型
            top3 = counter.most_common(3)
            main_products[line] = [p[0] for p in top3]
    return main_products


def analyze_production_change(folder: str, target_file: str, days: int = 40):
    all_files = get_ppt_files(folder, days)
    if not all_files:
        return None, "未找到有效PPT文件"
    main_products = calculate_main_products(all_files)
    target_path = os.path.join(folder, target_file)
    today_data = parse_ppt(target_path)
    results = []
    for line in sorted(today_data.keys()):
        products = today_data[line]
        main_list = main_products.get(line, [])  # 现在是前3名主力机型列表
        # 判断：当日生产的任一机型不在前3名主力中 → 需要重点检查
        is_change = any(p not in main_list for p in products)
        main_display = ' / '.join(main_list) if main_list else '未找到主力产品'
        results.append({
            'line': line,
            'main_product': main_display,
            'today_products': ' / '.join(products),
            'is_change': is_change,
            'status': '⚠️ 产品变化' if is_change else '✓ 正常'
        })
    return results, None


# ==================== PDF导出 ====================

def export_to_pdf(results, target_file, output_path=None):
    """导出分析结果到PDF"""
    if output_path is None:
        date_str = datetime.now().strftime('%Y%m%d')
        output_path = f"总组立课重点检查线体_{date_str}.pdf"
    
    # 尝试注册中文字体
    chinese_font = None
    font_paths = [
        '/usr/share/fonts/truetype/arphic/uming.ttc',
        '/usr/share/fonts/truetype/arphic/ukai.ttc',
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont('ChineseFont', fp))
                chinese_font = 'ChineseFont'
                break
            except:
                pass
    
    if chinese_font is None:
        # 没有中文字体，用内置字体
        chinese_font = 'Helvetica'
    
    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=15*mm,
        leftMargin=15*mm,
        topMargin=15*mm,
        bottomMargin=15*mm
    )
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontName=chinese_font,
        fontSize=16,
        alignment=1,  # 居中
        spaceAfter=10*mm
    )
    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Normal'],
        fontName=chinese_font,
        fontSize=10,
        alignment=1,
        spaceAfter=8*mm
    )
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontName=chinese_font,
        fontSize=9
    )
    
    elements = []
    
    # 标题
    date_match = re.search(r'(\d+)月(\d+)日', target_file)
    if date_match:
        report_date = f"2026年{date_match.group(1)}月{date_match.group(2)}日"
    else:
        report_date = datetime.now().strftime('%Y年%m月%d日')
    
    elements.append(Paragraph("总组立课自工程检查重点检查线体", title_style))
    elements.append(Paragraph(f"📅 报告日期：{report_date}  &nbsp;&nbsp;  📊 线体总数：{len(results)}  &nbsp;&nbsp;  ⚠️ 产品变化：{sum(1 for r in results if r['is_change'])}条", subtitle_style))
    elements.append(Spacer(1, 5*mm))
    
    # 统计摘要
    change_lines = [r for r in results if r['is_change']]
    normal_lines = [r for r in results if not r['is_change']]
    
    summary_text = f"<b>📋 分析摘要：</b> 共{len(results)}条线体，<font color='red'><b>{len(change_lines)}条</b></font>产品发生变化，{len(normal_lines)}条正常。"
    elements.append(Paragraph(summary_text, body_style))
    elements.append(Spacer(1, 5*mm))
    
    # 表格数据
    table_data = [['线体', '主力产品', '当日生产产品', '状态']]
    for r in results:
        status = '⚠️ 产品变化' if r['is_change'] else '✓ 正常'
        table_data.append([
            r['line'],
            r['main_product'],
            r['today_products'],
            status
        ])
    
    # 创建表格
    col_widths = [40*mm, 40*mm, 55*mm, 35*mm]
    table = Table(table_data, colWidths=col_widths)
    
    # 表格样式
    style_commands = [
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4682B4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), chinese_font),
        ('FONTNAME', (0, 1), (-1, -1), chinese_font),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('FONTSIZE', (0, 1), (-1, -1), 9),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
    ]
    
    # 重点检查行标红
    for i, r in enumerate(results, start=1):
        if r['is_change']:
            style_commands.append(('BACKGROUND', (0, i), (-1, i), colors.HexColor('#FFE4E1')))
            style_commands.append(('TEXTCOLOR', (3, i), (3, i), colors.HexColor('#B22222')))
    
    table.setStyle(TableStyle(style_commands))
    elements.append(table)
    elements.append(Spacer(1, 8*mm))
    
    # 变化点详情
    if change_lines:
        detail_text = "<b>⚠️ 产品变化线体详情：</b><br/>"
        for r in change_lines:
            detail_text += f"• {r['line']}：主力产品<b>{r['main_product']}</b> → 当日生产<b>{r['today_products']}</b><br/>"
        elements.append(Paragraph(detail_text, body_style))
    else:
        elements.append(Paragraph("<b>✓ 所有线体产品与主力一致，无产品变化。</b>", body_style))
    
    elements.append(Spacer(1, 5*mm))
    
    # 页脚
    footer_style = ParagraphStyle(
        'Footer',
        parent=styles['Normal'],
        fontName=chinese_font,
        fontSize=8,
        textColor=colors.grey,
        alignment=1
    )
    footer_text = f"本报告由总组立课自工程检查工具自动生成  |  生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    elements.append(Paragraph(footer_text, footer_style))
    
    doc.build(elements)
    return output_path


# ==================== 发送邮件 ====================

def send_email_with_pdf(pdf_path, recipient, subject=None, body=None):
    """发送PDF报告到邮箱"""
    config = load_email_config()
    
    if not config.get('password'):
        raise Exception("邮件配置不完整，请先设置发件邮箱和授权码")
    if not recipient:
        raise Exception("收件人邮箱不能为空")
    
    if subject is None:
        date_str = datetime.now().strftime('%Y年%m月%d日')
        subject = f"总组立课自工程检查重点检查线体 {date_str}"
    
    if body is None:
        body = f"您好，<br/>请查收附件中的总组立课自工程检查重点检查线体报告。<br/><br/>本邮件由自工程检查工具自动发送。"
    
    # 构建邮件
    msg = MIMEMultipart()
    msg['From'] = config['sender']
    msg['To'] = recipient
    msg['Subject'] = subject
    
    # HTML正文
    msg.attach(MIMEText(body, 'html', 'utf-8'))
    
    # 附件
    if os.path.exists(pdf_path):
        with open(pdf_path, 'rb') as f:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(f.read())
            encoders.encode_base64(part)
            filename = os.path.basename(pdf_path)
            part.add_header('Content-Disposition', f'attachment; filename="{filename}"')
            msg.attach(part)
    
    # 发送
    try:
        with smtplib.SMTP_SSL(config['smtp_server'], config['smtp_port']) as server:
            server.login(config['sender'], config['password'])
            server.sendmail(config['sender'], recipient, msg.as_string())
        return True
    except Exception as e:
        raise Exception(f"发送失败: {str(e)}")


# ==================== PyQt界面 ====================

COLORS = {
    'header': QColor(70, 130, 180),
    'normal_row': QColor(255, 255, 255),
    'change_row': QColor(255, 228, 225),
    'change_text': QColor(180, 34, 34),
}


class AnalyzeThread(QThread):
    finished = pyqtSignal(object)
    error = pyqtSignal(str)
    
    def __init__(self, folder, target_file, days=40):
        super().__init__()
        self.folder = folder
        self.target_file = target_file
        self.days = days
    
    def run(self):
        try:
            results, error = analyze_production_change(self.folder, self.target_file, self.days)
            if error:
                self.error.emit(error)
            else:
                self.finished.emit(results)
        except Exception as e:
            self.error.emit(str(e))


class EmailConfigDialog(QDialog):
    """邮件配置对话框"""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("📧 邮件配置")
        self.setFixedSize(450, 300)
        
        layout = QFormLayout(self)
        
        # 发件邮箱
        self.sender_input = QLineEdit()
        self.sender_input.setPlaceholderText("16835577@qq.com")
        layout.addRow("发件邮箱:", self.sender_input)
        
        # 授权码
        self.password_input = QLineEdit()
        self.password_input.setEchoMode(QLineEdit.Password)
        self.password_input.setPlaceholderText("QQ邮箱授权码（非登录密码）")
        layout.addRow("授权码:", self.password_input)
        
        # 收件邮箱
        self.recipient_input = QLineEdit()
        self.recipient_input.setPlaceholderText("收件人邮箱地址")
        layout.addRow("收件邮箱:", self.recipient_input)
        
        # 说明
        help_label = QLabel("💡 <a href='https://service.mail.qq.com/cgi-bin/help?subtype=1&&id=28&&no=1001256'>如何获取QQ邮箱授权码？</a>")
        help_label.setOpenExternalLinks(True)
        layout.addRow("", help_label)
        
        # 按钮
        buttons = QDialogButtonBox(
            QDialogButtonBox.Save | QDialogButtonBox.Cancel
        )
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addRow(buttons)
        
        # 加载已有配置
        self.load_config()
    
    def load_config(self):
        config = load_email_config()
        self.sender_input.setText(config.get('sender', ''))
        self.password_input.setText(config.get('password', ''))
        self.recipient_input.setText(config.get('recipient', ''))
    
    def save_config(self):
        config = {
            'sender': self.sender_input.text().strip(),
            'password': self.password_input.text().strip(),
            'recipient': self.recipient_input.text().strip(),
            'smtp_server': 'smtp.qq.com',
            'smtp_port': 465
        }
        save_email_config(config)


class SelfCheckWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.folder = ""
        self.ppt_files = []
        self.current_results = None
        self.current_pdf_path = None
        self.init_ui()
    
    def init_ui(self):
        self.setWindowTitle("总组立课自工程检查重点检查线体工具 v2.0")
        self.setGeometry(100, 100, 950, 700)
        
        central = QWidget()
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)
        
        # ---- 顶部：文件夹选择 ----
        top_layout = QHBoxLayout()
        self.folder_label = QLabel("未选择文件夹")
        self.folder_label.setStyleSheet("color: #666;")
        top_layout.addWidget(QLabel("📁 数据文件夹:"))
        top_layout.addWidget(self.folder_label, 1)
        self.select_folder_btn = QPushButton("选择文件夹")
        self.select_folder_btn.clicked.connect(self.select_folder)
        top_layout.addWidget(self.select_folder_btn)
        layout.addLayout(top_layout)
        
        # ---- 中间：控制按钮 ----
        ctrl_layout = QHBoxLayout()
        
        self.file_list_label = QLabel("请先选择文件夹")
        ctrl_layout.addWidget(QLabel("📋 选择文件:"))
        ctrl_layout.addWidget(self.file_list_label, 1)
        
        self.analyze_btn = QPushButton("🔍 分析变化点")
        self.analyze_btn.clicked.connect(self.run_analysis)
        self.analyze_btn.setEnabled(False)
        ctrl_layout.addWidget(self.analyze_btn)
        
        self.reload_btn = QPushButton("🔄 刷新")
        self.reload_btn.clicked.connect(self.reload_files)
        self.reload_btn.setEnabled(False)
        ctrl_layout.addWidget(self.reload_btn)
        
        layout.addLayout(ctrl_layout)
        
        # ---- 文件选择按钮 ----
        self.file_btn = QPushButton("=== 请选择文件 ===")
        self.file_btn.setEnabled(False)
        self.file_btn.clicked.connect(self.show_file_menu)
        layout.addWidget(self.file_btn)
        
        # ---- 结果表格 ----
        self.table = QTableWidget()
        self.table.setColumnCount(4)
        self.table.setHorizontalHeaderLabels(['线体', '主力产品', '当日生产产品', '状态'])
        self.table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        self.table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.table.setSelectionBehavior(QAbstractItemView.SelectRow)
        self.table.setAlternatingRowColors(True)
        layout.addWidget(self.table)
        
        # ---- 底部操作按钮 ----
        bottom_layout = QHBoxLayout()
        
        self.export_pdf_btn = QPushButton("📄 导出PDF")
        self.export_pdf_btn.clicked.connect(self.export_pdf)
        self.export_pdf_btn.setEnabled(False)
        bottom_layout.addWidget(self.export_pdf_btn)
        
        self.send_email_btn = QPushButton("📧 发送邮件")
        self.send_email_btn.clicked.connect(self.send_email)
        self.send_email_btn.setEnabled(False)
        bottom_layout.addWidget(self.send_email_btn)
        
        self.config_email_btn = QPushButton("⚙️ 邮件配置")
        self.config_email_btn.clicked.connect(self.config_email)
        bottom_layout.addWidget(self.config_email_btn)
        
        bottom_layout.addStretch()
        
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("就绪")
    
    def select_folder(self):
        folder = QFileDialog.getExistingDirectory(
            self, "选择生产计划文件夹", 
            os.path.expanduser("~/下载")
        )
        if folder:
            self.folder = folder
            self.folder_label.setText(folder)
            self.folder_label.setStyleSheet("color: #000;")
            self.reload_files()
    
    def reload_files(self):
        if not self.folder:
            return
        self.ppt_files = get_ppt_files(self.folder, days=40)
        if not self.ppt_files:
            self.file_list_label.setText("未找到PPT文件（检查是否在40天范围内）")
            self.file_btn.setEnabled(False)
            self.reload_btn.setEnabled(False)
            self.analyze_btn.setEnabled(False)
            return
        recent = self.ppt_files[-5:] if len(self.ppt_files) > 5 else self.ppt_files
        file_names = [f[2] for f in recent]
        self.file_list_label.setText(f"共 {len(self.ppt_files)} 个文件")
        self.file_btn.setEnabled(True)
        self.reload_btn.setEnabled(True)
        self.analyze_btn.setEnabled(False)
        if self.ppt_files:
            latest = self.ppt_files[-1][2]
            self.file_btn.setText(f"📄 {latest} ▼")
            self.selected_file = latest
            self.analyze_btn.setEnabled(True)
        self.status_bar.showMessage(f"已加载 {len(self.ppt_files)} 个文件")
    
    def show_file_menu(self):
        if not self.ppt_files:
            return
        from PyQt5.QtWidgets import QMenu
        menu = QMenu(self)
        for date, filepath, filename in reversed(self.ppt_files):
            action = menu.addAction(filename)
            action.triggered.connect(lambda checked, f=filename: self.select_file(f))
        menu.exec_(self.file_btn.mapToGlobal(self.file_btn.rect().bottomLeft()))
    
    def select_file(self, filename):
        self.selected_file = filename
        self.file_btn.setText(f"📄 {filename}")
        self.analyze_btn.setEnabled(True)
        self.status_bar.showMessage(f"已选择: {filename}")
    
    def run_analysis(self):
        if not self.folder or not hasattr(self, 'selected_file'):
            return
        self.analyze_btn.setEnabled(False)
        self.status_bar.showMessage("分析中...")
        progress = QProgressDialog("正在分析...", "取消", 0, 0, self)
        progress.setWindowTitle("处理中")
        progress.setWindowModality(Qt.WindowModal)
        progress.show()
        self.thread = AnalyzeThread(self.folder, self.selected_file, days=40)
        self.thread.finished.connect(lambda r: self.on_analysis_done(r, progress))
        self.thread.error.connect(lambda e: self.on_analysis_error(e, progress))
        self.thread.start()
    
    def on_analysis_done(self, results, progress):
        progress.close()
        self.analyze_btn.setEnabled(True)
        if not results:
            self.status_bar.showMessage("分析完成，无数据")
            return
        self.current_results = results
        self.display_results(results)
        changes = sum(1 for r in results if r['is_change'])
        self.status_bar.showMessage(f"✅ 分析完成！{len(results)}条线体，{changes}个变化点")
        self.export_pdf_btn.setEnabled(True)
        self.send_email_btn.setEnabled(True)
    
    def on_analysis_error(self, error, progress):
        progress.close()
        self.analyze_btn.setEnabled(True)
        QMessageBox.warning(self, "错误", f"分析失败:\n{error}")
        self.status_bar.showMessage("分析失败")
    
    def display_results(self, results):
        self.table.setRowCount(len(results))
        for i, row_data in enumerate(results):
            items = [
                row_data['line'],
                row_data['main_product'],
                row_data['today_products'],
                row_data['status']
            ]
            for j, text in enumerate(items):
                item = QTableWidgetItem(text)
                item.setTextAlignment(Qt.AlignCenter)
                if row_data['is_change']:
                    item.setBackground(COLORS['change_row'])
                    if j == 3:
                        item.setForeground(COLORS['change_text'])
                        font = QFont()
                        font.setBold(True)
                        item.setFont(font)
                else:
                    item.setBackground(COLORS['normal_row'])
                self.table.setItem(i, j, item)
        self.table.resizeColumnsToContents()
    
    def export_pdf(self):
        if not self.current_results:
            return
        # 选择保存位置
        date_str = datetime.now().strftime('%Y%m%d')
        default_name = f"总组立课重点检查线体_{date_str}.pdf"
        path, _ = QFileDialog.getSaveFileName(
            self, "保存PDF报告", default_name, "PDF Files (*.pdf)"
        )
        if not path:
            return
        try:
            self.status_bar.showMessage("正在生成PDF...")
            output_path = export_to_pdf(self.current_results, self.selected_file, path)
            self.current_pdf_path = output_path
            self.status_bar.showMessage(f"✅ PDF已保存: {output_path}")
            # 自动打开PDF
            try:
                if sys.platform == 'darwin':
                    subprocess.run(['open', output_path])
                elif sys.platform == 'win32':
                    os.startfile(output_path)
                else:
                    subprocess.run(['xdg-open', output_path])
            except Exception:
                pass  # 打开失败不影响主流程
            QMessageBox.information(self, "成功", f"PDF报告已保存:\n{output_path}")
        except Exception as e:
            QMessageBox.warning(self, "错误", f"导出PDF失败:\n{str(e)}")
            self.status_bar.showMessage("PDF导出失败")
    
    def config_email(self):
        dialog = EmailConfigDialog(self)
        if dialog.exec_():
            dialog.save_config()
            QMessageBox.information(self, "成功", "邮件配置已保存")
    
    def send_email(self):
        if not self.current_results:
            return
        config = load_email_config()
        if not config.get('password'):
            QMessageBox.warning(self, "请先配置", "请先设置邮件配置（发件邮箱和授权码）")
            self.config_email()
            return
        if not config.get('recipient'):
            QMessageBox.warning(self, "请先配置", "请先设置收件人邮箱")
            self.config_email()
            return
        # 先生成PDF
        try:
            self.status_bar.showMessage("正在生成PDF...")
            date_str = datetime.now().strftime('%Y%m%d')
            temp_pdf = f"/tmp/总组立课重点检查线体_{date_str}.pdf"
            output_path = export_to_pdf(self.current_results, self.selected_file, temp_pdf)
            self.current_pdf_path = output_path
        except Exception as e:
            QMessageBox.warning(self, "错误", f"生成PDF失败:\n{str(e)}")
            return
        # 发送邮件
        try:
            self.status_bar.showMessage("正在发送邮件...")
            send_email_with_pdf(
                self.current_pdf_path,
                config['recipient']
            )
            self.status_bar.showMessage(f"✅ 邮件已发送至: {config['recipient']}")
            QMessageBox.information(self, "成功", f"邮件已发送至:\n{config['recipient']}")
        except Exception as e:
            QMessageBox.warning(self, "错误", f"发送失败:\n{str(e)}")
            self.status_bar.showMessage("邮件发送失败")


def main():
    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    font = QFont()
    font.setFamily('Microsoft YaHei UI')
    app.setFont(font)
    window = SelfCheckWindow()
    window.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    main()
