#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
自工程检查重点检查线体工具 v3.1
- 读取生产计划PPT
- 根据过去30天数据计算各线体主力产品
- 标记产品变化线体为重点检查对象
- 琴型标准化处理：35B=NP-35B，显示完整名称
- 读取琴型关注点表，自动匹配关注点
- 支持导出PDF、发送邮件
"""

import os
import sys
import re
import smtplib
import json
import subprocess
from datetime import datetime, timedelta
from collections import defaultdict, Counter
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLabel, QTableWidget, QTableWidgetItem, QFileDialog,
    QMessageBox, QHeaderView, QAbstractItemView, QStatusBar, QProgressDialog,
    QDialog, QFormLayout, QLineEdit, QDialogButtonBox
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal
from PyQt5.QtGui import QColor, QFont

from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


# ==================== 琴型标准化处理 ====================

def normalize_product_name(product: str) -> str:
    """
    标准化琴型名称，用于统计匹配
    例如：NP-35B -> 35B，NP35B -> 35B
    规则：去除系列前缀和连接符，只保留核心型号
    """
    if not product:
        return product
    product = product.strip()
    # 去除常见前缀的连接符
    patterns = [
        r'^P[- ]*',
        r'^PK[- ]*',
        r'^CL[- ]*',
        r'^NP[- ]*',
        r'^PSR[- ]*',
        r'^KBP[- ]*',
        r'^YPT[- ]*',
        r'^EZ[- ]*',
        r'^KB[- ]*',
    ]
    normalized = product
    for pattern in patterns:
        normalized = re.sub(pattern, '', normalized, flags=re.IGNORECASE)
    return normalized.strip()


# ==================== 琴型关注点配置 ====================
ATTENTION_FILE = os.path.expanduser("~/.ipqa_attention_points.xlsx")


def load_attention_points():
    """加载琴型关注点表"""
    if not os.path.exists(ATTENTION_FILE):
        return {}
    try:
        import pandas as pd
        df = pd.read_excel(ATTENTION_FILE, sheet_name='按系列统计', header=None)
        attention_map = {}
        for idx, row in df.iterrows():
            if idx < 3:
                continue
            series = str(row[0]).strip() if pd.notna(row[0]) else ""
            product = str(row[1]).strip() if pd.notna(row[1]) else ""
            attention = str(row[2]).strip() if pd.notna(row[2]) and row[2] != "关注点" else ""
            if product and attention and product != "琴型":
                norm = normalize_product_name(product)
                attention_map[norm] = attention
                attention_map[product] = attention
                attention_map[f"{series}|{product}"] = attention
        return attention_map
    except Exception as e:
        print(f"读取关注点表失败: {e}")
        return {}


def get_attention_point(series, product, attention_map):
    """获取关注点"""
    norm_product = normalize_product_name(product)
    key1 = f"{series}|{product}"
    if key1 in attention_map:
        return attention_map[key1]
    key2 = f"{series}|{norm_product}"
    if key2 in attention_map:
        return attention_map[key2]
    if product in attention_map:
        return attention_map[product]
    if norm_product in attention_map:
        return attention_map[norm_product]
    for key, value in attention_map.items():
        if '|' in key:
            _, p = key.split('|', 1)
            if normalize_product_name(p) == norm_product:
                return value
        else:
            if normalize_product_name(key) == norm_product:
                return value
    return ""


# ==================== 邮件配置 ====================
EMAIL_CONFIG_FILE = os.path.expanduser("~/.ipqa_email_config.json")
ATTENTION_CONFIG_FILE = os.path.expanduser("~/.ipqa_attention_config.json")
DATA_FOLDER_CONFIG_FILE = os.path.expanduser("~/.ipqa_data_folder_config.json")


def load_email_config():
    if os.path.exists(EMAIL_CONFIG_FILE):
        with open(EMAIL_CONFIG_FILE, 'r') as f:
            return json.load(f)
    return {"smtp_server": "smtp.qq.com", "smtp_port": 465, "sender": "16835577@qq.com", "password": "", "recipient": ""}


def save_email_config(config):
    with open(EMAIL_CONFIG_FILE, 'w') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


def load_data_folder_config():
    if os.path.exists(DATA_FOLDER_CONFIG_FILE):
        with open(DATA_FOLDER_CONFIG_FILE, 'r') as f:
            return json.load(f)
    return {"data_folder": ""}


def save_data_folder_config(config):
    with open(DATA_FOLDER_CONFIG_FILE, 'w') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


def load_attention_config():
    if os.path.exists(ATTENTION_CONFIG_FILE):
        with open(ATTENTION_CONFIG_FILE, 'r') as f:
            return json.load(f)
    return {"attention_file": ""}


def save_attention_config(config):
    with open(ATTENTION_CONFIG_FILE, 'w') as f:
        json.dump(config, f, ensure_ascii=False, indent=2)


# ==================== PPT解析逻辑 ====================
EXCLUDE_KEYWORDS = ['社内', '亏产', '艾航', '中兴', '永信', '余裕人员', '培训', '休']
EXCLUDE_LINE_KEYWORDS = ['键盘', '踏板']
NON_PRODUCTION_KEYWORDS = ['培训', '休', '待机', '整理', '早会', '例会', '点检', '换线', '调试']


def is_valid_product(product: str) -> bool:
    if not product or not product.strip():
        return False
    product = product.strip()
    for kw in EXCLUDE_KEYWORDS:
        if kw in product:
            return False
    if re.match(r'^[\d\-\/]+$', product):
        return False
    return True


def is_production_line(line: str) -> bool:
    if not line or not line.strip():
        return False
    line = line.strip()
    for kw in NON_PRODUCTION_KEYWORDS:
        if kw in line:
            return False
    return True


def is_valid_line(line: str) -> bool:
    if not line or not line.strip():
        return False
    for kw in EXCLUDE_LINE_KEYWORDS:
        if kw in line:
            return False
    return True


def parse_ppt(filepath: str) -> dict:
    result = {}
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = list(table.rows)
                    if len(rows) < 2:
                        continue
                    header = [cell.text.strip() for cell in rows[0].cells]
                    try:
                        line_idx = header.index('线体')
                        product_idx = header.index('当日生产机种')
                    except ValueError:
                        continue
                    for row in rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if len(cells) <= max(line_idx, product_idx):
                            continue
                        line = cells[line_idx]
                        product = cells[product_idx]
                        if is_valid_line(line) and is_valid_product(product):
                            products = [p.strip() for p in product.split('/') if is_valid_product(p.strip())]
                            if products:
                                result[line] = products
    except Exception as e:
        print(f"解析PPT出错 {filepath}: {e}")
    return result


def extract_date_from_filename(filename: str) -> datetime:
    match = re.search(r'(\d+)月(\d+)日', filename)
    if match:
        return datetime(2026, int(match.group(1)), int(match.group(2)))
    return datetime(2026, 1, 1)


def get_ppt_files(folder: str, days: int = 30) -> list:
    files = []
    if not os.path.exists(folder):
        return files
    today = datetime.now()
    cutoff = today - timedelta(days=days)
    for f in os.listdir(folder):
        if f.endswith('.pptx'):
            filepath = os.path.join(folder, f)
            file_date = extract_date_from_filename(f)
            if file_date >= cutoff:
                files.append((file_date, filepath, f))
    files.sort(key=lambda x: x[0])
    return files


def calculate_main_products(ppt_files: list) -> dict:
    """计算每条线体的前3名主力机型（标准化后合并统计）"""
    # 原始产品名记录
    line_product_original = defaultdict(Counter)
    
    for _, filepath, _ in ppt_files:
        daily_data = parse_ppt(filepath)
        for line, products in daily_data.items():
            if not is_production_line(line):
                continue
            for product in products:
                line_product_original[line][product] += 1
    
    # 合并标准化名称的计数，并保留显示名（优先选最长/最完整的名称）
    main_products = {}
    for line, counter in line_product_original.items():
        merged = Counter()
        orig_of_norm = {}  # norm -> longest original name
        for product, count in counter.items():
            norm = normalize_product_name(product)
            merged[norm] += count
            if norm not in orig_of_norm or len(product) > len(orig_of_norm[norm]):
                orig_of_norm[norm] = product
        
        # 过滤掉生产次数小于5次的产品，按生产次数排序后取前3
        filtered = [(norm, count) for norm, count in merged.items() if count >= 5]
        filtered.sort(key=lambda x: x[1], reverse=True)
        top3 = filtered[:3] if filtered else []
        main_products[line] = [orig_of_norm.get(n[0], n[0]) for n in top3]
    
    return main_products


def analyze_production_change(folder: str, target_file: str, days: int = 40, attention_map: dict = None):
    if attention_map is None:
        attention_map = {}
    
    all_files = get_ppt_files(folder, days)
    if not all_files:
        return None, "未找到有效PPT文件"
    
    main_products = calculate_main_products(all_files)
    target_path = os.path.join(folder, target_file)
    today_data = parse_ppt(target_path)
    results = []
    
    for line in sorted(today_data.keys()):
        products = today_data[line]
        main_list = main_products.get(line, [])
        
        # 标准化比较，判断是否变化
        is_change = False
        for p in products:
            p_norm = normalize_product_name(p)
            found = any(normalize_product_name(m) == p_norm for m in main_list)
            if not found:
                is_change = True
                break
        
        main_display = ' / '.join(main_list) if main_list else '未找到主力产品'
        
        # 获取关注点
        attention_points = []
        for p in products:
            att = get_attention_point(line, p, attention_map)
            if att and att not in attention_points:
                attention_points.append(att)
        attention_display = ' '.join(attention_points) if attention_points else ''
        
        results.append({
            'line': line,
            'main_product': main_display,
            'today_products': ' / '.join(products),
            'is_change': is_change,
            'status': '⚠️ 产品变化' if is_change else '✓ 正常',
            'attention': attention_display
        })
    
    return results, None


# ==================== PDF导出 ====================
def export_to_pdf(results, target_file, output_path=None):
    if output_path is None:
        date_str = datetime.now().strftime('%Y%m%d')
        output_path = f"总组立课重点检查线体_{date_str}.pdf"
    
    chinese_font = None
    font_paths = [
        'C:/Windows/Fonts/simhei.ttf',
        'C:/Windows/Fonts/msyh.ttc',
        'C:/Windows/Fonts/simsun.ttc',
        '/usr/share/fonts/truetype/arphic/uming.ttc',
        '/usr/share/fonts/truetype/arphic/ukai.ttc',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc'
    ]
    for fp in font_paths:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont('ChineseFont', fp))
                chinese_font = 'ChineseFont'
                break
            except:
                pass
    if chinese_font is None:
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            pdfmetrics.registerFont(TTFont('ChineseFont', 'C:/Windows/Fonts/simhei.ttf'))
            chinese_font = 'ChineseFont'
        except:
            chinese_font = 'Helvetica'
    
    doc = SimpleDocTemplate(output_path, pagesize=landscape(A4), rightMargin=12*mm, leftMargin=12*mm, topMargin=12*mm, bottomMargin=12*mm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontName=chinese_font, fontSize=16, alignment=1, spaceAfter=8*mm)
    subtitle_style = ParagraphStyle('CustomSubtitle', parent=styles['Normal'], fontName=chinese_font, fontSize=10, alignment=1, spaceAfter=6*mm)
    body_style = ParagraphStyle('CustomBody', parent=styles['Normal'], fontName=chinese_font, fontSize=9)
    
    elements = []
    date_match = re.search(r'(\d+)月(\d+)日', target_file)
    report_date = f"2026年{date_match.group(1)}月{date_match.group(2)}日" if date_match else datetime.now().strftime('%Y年%m月%d日')
    
    elements.append(Paragraph("总组立课自工程检查重点检查线体", title_style))
    elements.append(Paragraph(f"报告日期：{report_date}    线体总数：{len(results)}    产品变化：{sum(1 for r in results if r['is_change'])}条", subtitle_style))
    elements.append(Spacer(1, 4*mm))
    
    change_lines = [r for r in results if r['is_change']]
    normal_lines = [r for r in results if not r['is_change']]
    elements.append(Paragraph(f"<b>分析摘要：</b> 共{len(results)}条线体，<font color='red'><b>{len(change_lines)}条</b></font>产品发生变化，{len(normal_lines)}条正常。", body_style))
    elements.append(Spacer(1, 4*mm))
    
    header_style = ParagraphStyle('Header', fontName=chinese_font, fontSize=9, leading=11, alignment=1)
    table_data = [['线体', '主力产品', '当日生产产品', '状态', '关注点']]
    table_data[0] = [Paragraph(h, header_style) for h in table_data[0]]
    for r in results:
        status = '产品变化' if r['is_change'] else '无异常'
        if r['is_change'] and r['attention']:
            attention_text = r['attention'][:50] + '...' if len(r['attention']) > 50 else r['attention']
        else:
            attention_text = '/'
        cell_style = ParagraphStyle('Cell', fontName=chinese_font, fontSize=7, leading=9, alignment=1)
        table_data.append([
            Paragraph(r['line'], cell_style),
            Paragraph(r['main_product'], cell_style),
            Paragraph(r['today_products'], cell_style),
            Paragraph(status, cell_style),
            Paragraph(attention_text, cell_style)
        ])
    
    col_widths = [35*mm, 50*mm, 85*mm, 28*mm, 64*mm]
    table = Table(table_data, colWidths=col_widths)
    style_commands = [
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4682B4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
        ('FONTNAME', (0, 0), (-1, 0), chinese_font),
        ('FONTNAME', (0, 1), (-1, -1), chinese_font),
        ('FONTSIZE', (0, 0), (-1, 0), 9),
        ('FONTSIZE', (0, 1), (-1, -1), 7),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('TOPPADDING', (0, 0), (-1, -1), 2),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
        ('LEFTPADDING', (0, 0), (-1, -1), 1),
        ('RIGHTPADDING', (0, 0), (-1, -1), 1),
    ]
    for i, r in enumerate(results, start=1):
        if r['is_change']:
            style_commands.append(('BACKGROUND', (0, i), (-1, i), colors.HexColor('#FFFFCC')))
            style_commands.append(('TEXTCOLOR', (3, i), (3, i), colors.HexColor('#CC0000')))
    table.setStyle(TableStyle(style_commands))
    elements.append(table)
    elements.append(Spacer(1, 5*mm))
    
    if change_lines:
        detail_text = "<b>产品变化线体详情：</b><br/>"
        for r in change_lines:
            att_info = f"<br/>&nbsp;&nbsp;&nbsp;&nbsp;<font color='blue'>关注点：{r['attention']}</font>" if r['attention'] else ""
            detail_text += f"• {r['line']}：主力 <b>{r['main_product']}</b> → 当日 <b>{r['today_products']}</b>{att_info}<br/>"
        elements.append(Paragraph(detail_text, body_style))
    else:
        elements.append(Paragraph("<b>所有线体产品与主力一致，无产品变化。</b>", body_style))
    
    elements.append(Spacer(1, 4*mm))
    footer_style = ParagraphStyle('Footer', parent=styles['Normal'], fontName=chinese_font, fontSize=8, textColor=colors.grey, alignment=1)
    elements.append(Paragraph(f"本报告由总组立课自工程检查工具自动生成  |  生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", footer_style))
    
    doc.build(elements)
    return output_path


# ==================== 发送邮件 ====================
def send_email_with_pdf(pdf_path, recipient, subject=None, body=None):
    config = load_email_config()
    if not config.get('password'):
        raise Exception("邮件配置不完整")
    if not recipient:
        raise Exception("收件人邮箱不能为空")
    if subject is None:
        date_str = datetime.now().strftime('%Y年%m月%d日')
        subject = f"总组立课自工程检查重点检查线体 {date_str}"
    if body is None:
        body = "您好，<br/>请查收附件中的总组立课自工程检查重点检查线体报告。<br/><br/>本邮件由自工程检查工具自动发送。"
    
    msg = MIMEMultipart()
    msg['From'] = config['sender']
    msg['To'] = recipient
    msg['Subject'] = subject
    msg.attach(MIMEText(body, 'html', 'utf-8'))
    
    if os.path.exists(pdf_path):
        with open(pdf_path, 'rb') as f:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(f.read())
            encoders.encode_base64(part)
            part.add_header('Content-Disposition', f'attachment; filename="{os.path.basename(pdf_path)}"')
            msg.attach(part)
    
    with smtplib.SMTP_SSL(config['smtp_server'], config['smtp_port']) as server:
        server.login(config['sender'], config['password'])
        server.sendmail(config['sender'], recipient, msg.as_string())
    return True


# ==================== PyQt界面 ====================
COLORS = {'header': QColor(70, 130, 180), 'normal_row': QColor(255, 255, 255), 'change_row': QColor(255, 228, 225), 'change_text': QColor(180, 34, 34)}


class AnalyzeThread(QThread):
    finished = pyqtSignal(object)
    error = pyqtSignal(str)
    
    def __init__(self, folder, target_file, days=30, attention_map=None):
        super().__init__()
        self.folder = folder
        self.target_file = target_file
        self.days = days
        self.attention_map = attention_map or {}
    
    def run(self):
        try:
            results, error = analyze_production_change(self.folder, self.target_file, self.days, self.attention_map)
            if error:
                self.error.emit(error)
            else:
                self.finished.emit(results)
        except Exception as e:
            self.error.emit(str(e))


class EmailConfigDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("📧 邮件配置")
        self.setFixedSize(450, 300)
        layout = QFormLayout(self)
        self.sender_input = QLineEdit()
        self.sender_input.setPlaceholderText("16835577@qq.com")
        layout.addRow("发件邮箱:", self.sender_input)
        self.password_input = QLineEdit()
        self.password_input.setEchoMode(QLineEdit.Password)
        self.password_input.setPlaceholderText("QQ邮箱授权码")
        layout.addRow("授权码:", self.password_input)
        self.recipient_input = QLineEdit()
        self.recipient_input.setPlaceholderText("收件人邮箱地址")
        layout.addRow("收件邮箱:", self.recipient_input)
        help_label = QLabel("💡 <a href='https://service.mail.qq.com/cgi-bin/help?subtype=1&&id=28&&no=1001256'>如何获取QQ邮箱授权码？</a>")
        help_label.setOpenExternalLinks(True)
        layout.addRow("", help_label)
        buttons = QDialogButtonBox(QDialogButtonBox.Save | QDialogButtonBox.Cancel)
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)
        self.load_config()
    
    def load_config(self):
        config = load_email_config()
        self.sender_input.setText(config.get('sender', ''))
        self.password_input.setText(config.get('password', ''))
        self.recipient_input.setText(config.get('recipient', ''))
    
    def save_config(self):
        config = {'sender': self.sender_input.text().strip(), 'password': self.password_input.text().strip(), 'recipient': self.recipient_input.text().strip(), 'smtp_server': 'smtp.qq.com', 'smtp_port': 465}
        save_email_config(config)


class AttentionFileDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("琴型关注点配置")
        self.setFixedSize(500, 150)
        layout = QVBoxLayout(self)
        layout.addWidget(QLabel("选择琴型关注点Excel文件："))
        path_layout = QHBoxLayout()
        self.path_label = QLabel("未选择文件")
        self.path_label.setStyleSheet("color: #666;")
        path_layout.addWidget(self.path_label, 1)
        self.browse_btn = QPushButton("浏览...")
        self.browse_btn.clicked.connect(self.browse)
        path_layout.addWidget(self.browse_btn)
        layout.addLayout(path_layout)
        layout.addWidget(QLabel("提示：文件应包含'系列'、'琴型'、'关注点'三列"))
        buttons = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)
        self.current_path = ""
        self.load_config()
    
    def load_config(self):
        config = load_attention_config()
        if config.get('attention_file') and os.path.exists(config['attention_file']):
            self.current_path = config['attention_file']
            self.path_label.setText(self.current_path)
            self.path_label.setStyleSheet("color: #000;")
    
    def save_config(self):
        save_attention_config({"attention_file": self.current_path})
    
    def browse(self):
        path, _ = QFileDialog.getOpenFileName(self, "选择琴型关注点文件", os.path.expanduser("~/下载"), "Excel Files (*.xlsx *.xls)")
        if path:
            self.current_path = path
            self.path_label.setText(path)
            self.path_label.setStyleSheet("color: #000;")


class SelfCheckWindow(QMainWindow):
    def __init__(self):
        try:
            print("SelfCheckWindow.__init__ starting...")
            super().__init__()
            print("Super class initialized")
            self.folder = ""
            self.ppt_files = []
            self.current_results = None
            self.current_pdf_path = None
            self.attention_map = {}
            print("Variables initialized")
            print("Calling init_ui...")
            self.init_ui()
            print("init_ui completed")
            print("Calling load_attention_file...")
            self.load_attention_file()
            print("load_attention_file completed")
            print("Calling load_saved_folder...")
            self.load_saved_folder()
            print("load_saved_folder completed")
            print("SelfCheckWindow.__init__ completed")
        except Exception as e:
            print(f"Error in SelfCheckWindow.__init__: {e}")
            import traceback
            traceback.print_exc()
            raise
    
    def load_saved_folder(self):
        config = load_data_folder_config()
        if config.get('data_folder') and os.path.exists(config['data_folder']):
            self.folder = config['data_folder']
            self.folder_label.setText(self.folder)
            self.folder_label.setStyleSheet("color: #000;")
            self.reload_files()
    
    def load_attention_file(self):
        global ATTENTION_FILE
        config = load_attention_config()
        if config.get('attention_file') and os.path.exists(config['attention_file']):
            ATTENTION_FILE = config['attention_file']
        if os.path.exists(ATTENTION_FILE):
            try:
                import pandas as pd
                df = pd.read_excel(ATTENTION_FILE, sheet_name='按系列统计', header=None)
                for idx, row in df.iterrows():
                    if idx < 3:
                        continue
                    series = str(row[0]).strip() if pd.notna(row[0]) else ""
                    product = str(row[1]).strip() if pd.notna(row[1]) else ""
                    attention = str(row[2]).strip() if pd.notna(row[2]) and row[2] != "关注点" else ""
                    if product and attention and product != "琴型":
                        norm = normalize_product_name(product)
                        self.attention_map[norm] = attention
                        self.attention_map[product] = attention
                        self.attention_map[f"{series}|{product}"] = attention
                if self.attention_map:
                    self.att_label.setText(f"关注点: 已加载 {len(self.attention_map)} 条")
                    self.att_label.setStyleSheet("color: #008000;")
                    self.status_bar.showMessage(f"已加载 {len(self.attention_map)} 条关注点")
                else:
                    self.att_label.setText("关注点: 文件为空")
                    self.att_label.setStyleSheet("color: #CC0000;")
                    self.status_bar.showMessage("关注点文件为空或格式错误")
            except Exception as e:
                error_msg = str(e)
                self.att_label.setText(f"关注点: 加载失败 - {error_msg[:30]}")
                self.att_label.setStyleSheet("color: #CC0000;")
                self.status_bar.showMessage(f"加载关注点失败: {error_msg}")
        else:
            self.att_label.setText("关注点: 未设置")
            self.att_label.setStyleSheet("color: #666;")
    
    def init_ui(self):
        self.setWindowTitle("总组立课自工程检查重点检查线体工具 v3.1")
        self.setGeometry(100, 100, 1100, 750)
        central = QWidget()
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)
        
        top_layout = QHBoxLayout()
        self.folder_label = QLabel("未选择文件夹")
        self.folder_label.setStyleSheet("color: #666;")
        top_layout.addWidget(QLabel("📁 数据文件夹:"))
        top_layout.addWidget(self.folder_label, 1)
        self.select_folder_btn = QPushButton("选择文件夹")
        self.select_folder_btn.clicked.connect(self.select_folder)
        top_layout.addWidget(self.select_folder_btn)
        layout.addLayout(top_layout)
        
        att_layout = QHBoxLayout()
        self.att_label = QLabel("📋 关注点: 未加载")
        self.att_label.setStyleSheet("color: #666;")
        att_layout.addWidget(self.att_label)
        self.att_btn = QPushButton("设置关注点文件")
        self.att_btn.clicked.connect(self.set_attention_file)
        att_layout.addWidget(self.att_btn)
        layout.addLayout(att_layout)
        
        ctrl_layout = QHBoxLayout()
        self.file_list_label = QLabel("请先选择文件夹")
        ctrl_layout.addWidget(QLabel("📋 选择文件:"))
        ctrl_layout.addWidget(self.file_list_label, 1)
        self.analyze_btn = QPushButton("🔍 分析变化点")
        self.analyze_btn.clicked.connect(self.run_analysis)
        self.analyze_btn.setEnabled(False)
        ctrl_layout.addWidget(self.analyze_btn)
        self.reload_btn = QPushButton("🔄 刷新")
        self.reload_btn.clicked.connect(self.reload_files)
        self.reload_btn.setEnabled(False)
        ctrl_layout.addWidget(self.reload_btn)
        layout.addLayout(ctrl_layout)
        
        self.file_btn = QPushButton("=== 请选择文件 ===")
        self.file_btn.setEnabled(False)
        self.file_btn.clicked.connect(self.show_file_menu)
        layout.addWidget(self.file_btn)
        
        self.table = QTableWidget()
        self.table.setColumnCount(5)
        self.table.setHorizontalHeaderLabels(['线体', '主力产品', '当日生产产品', '关注点', '状态'])
        self.table.horizontalHeader().setSectionResizeMode(QHeaderView.Stretch)
        self.table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.table.setAlternatingRowColors(True)
        layout.addWidget(self.table)
        
        bottom_layout = QHBoxLayout()
        self.export_pdf_btn = QPushButton("📄 导出PDF")
        self.export_pdf_btn.clicked.connect(self.export_pdf)
        self.export_pdf_btn.setEnabled(False)
        bottom_layout.addWidget(self.export_pdf_btn)
        self.send_email_btn = QPushButton("📧 发送邮件")
        self.send_email_btn.clicked.connect(self.send_email)
        self.send_email_btn.setEnabled(False)
        bottom_layout.addWidget(self.send_email_btn)
        self.config_email_btn = QPushButton("⚙️ 邮件配置")
        self.config_email_btn.clicked.connect(self.config_email)
        bottom_layout.addWidget(self.config_email_btn)
        bottom_layout.addStretch()
        layout.addLayout(bottom_layout)
        
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("就绪")
    
    def set_attention_file(self):
        dialog = AttentionFileDialog(self)
        if dialog.exec_():
            dialog.save_config()
            global ATTENTION_FILE
            ATTENTION_FILE = dialog.current_path
            self.attention_map = {}
            self.load_attention_file()
            QMessageBox.information(self, "成功", f"已设置关注点文件:\n{ATTENTION_FILE}")
    
    def select_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "选择生产计划文件夹", os.path.expanduser("~/下载"))
        if folder:
            self.folder = folder
            self.folder_label.setText(folder)
            self.folder_label.setStyleSheet("color: #000;")
            save_data_folder_config({"data_folder": folder})
            self.reload_files()
    
    def reload_files(self):
        if not self.folder:
            return
        self.ppt_files = get_ppt_files(self.folder)
        if not self.ppt_files:
            self.file_list_label.setText("未找到PPT文件")
            self.file_btn.setEnabled(False)
            self.reload_btn.setEnabled(False)
            self.analyze_btn.setEnabled(False)
            return
        self.file_list_label.setText(f"共 {len(self.ppt_files)} 个文件")
        self.file_btn.setEnabled(True)
        self.reload_btn.setEnabled(True)
        self.analyze_btn.setEnabled(False)
        if self.ppt_files:
            latest = self.ppt_files[-1][2]
            self.file_btn.setText(f"📄 {latest} ▼")
            self.selected_file = latest
            self.analyze_btn.setEnabled(True)
        self.status_bar.showMessage(f"已加载 {len(self.ppt_files)} 个文件")
    
    def show_file_menu(self):
        if not self.ppt_files:
            return
        from PyQt5.QtWidgets import QMenu
        menu = QMenu(self)
        for date, filepath, filename in reversed(self.ppt_files):
            action = menu.addAction(filename)
            action.triggered.connect(lambda checked, f=filename: self.select_file(f))
        menu.exec_(self.file_btn.mapToGlobal(self.file_btn.rect().bottomLeft()))
    
    def select_file(self, filename):
        self.selected_file = filename
        self.file_btn.setText(f"📄 {filename}")
        self.analyze_btn.setEnabled(True)
        self.status_bar.showMessage(f"已选择: {filename}")
    
    def run_analysis(self):
        if not self.folder or not hasattr(self, 'selected_file'):
            return
        self.analyze_btn.setEnabled(False)
        self.status_bar.showMessage("分析中...")
        progress = QProgressDialog("正在分析...", "取消", 0, 0, self)
        progress.setWindowTitle("处理中")
        progress.setWindowModality(Qt.WindowModal)
        progress.show()
        self.thread = AnalyzeThread(self.folder, self.selected_file, attention_map=self.attention_map)
        self.thread.finished.connect(lambda r: self.on_analysis_done(r, progress))
        self.thread.error.connect(lambda e: self.on_analysis_error(e, progress))
        self.thread.start()
    
    def on_analysis_done(self, results, progress):
        progress.close()
        self.analyze_btn.setEnabled(True)
        if not results:
            self.status_bar.showMessage("分析完成，无数据")
            return
        self.current_results = results
        self.display_results(results)
        changes = sum(1 for r in results if r['is_change'])
        self.status_bar.showMessage(f"✅ 分析完成！{len(results)}条线体，{changes}个变化点")
        self.export_pdf_btn.setEnabled(True)
        self.send_email_btn.setEnabled(True)
    
    def on_analysis_error(self, error, progress):
        progress.close()
        self.analyze_btn.setEnabled(True)
        QMessageBox.warning(self, "错误", f"分析失败:\n{error}")
        self.status_bar.showMessage("分析失败")
    
    def display_results(self, results):
        self.table.setRowCount(len(results))
        for i, row_data in enumerate(results):
            items = [row_data['line'], row_data['main_product'], row_data['today_products'],
                     row_data['attention'][:30] + '...' if len(row_data['attention']) > 30 else row_data['attention'],
                     row_data['status']]
            for j, text in enumerate(items):
                item = QTableWidgetItem(text)
                item.setTextAlignment(Qt.AlignCenter)
                if row_data['is_change']:
                    item.setBackground(COLORS['change_row'])
                    if j == 4:
                        item.setForeground(COLORS['change_text'])
                        font = QFont()
                        font.setBold(True)
                        item.setFont(font)
                else:
                    item.setBackground(COLORS['normal_row'])
                self.table.setItem(i, j, item)
        self.table.resizeColumnsToContents()
    
    def export_pdf(self):
        if not self.current_results:
            return
        date_str = datetime.now().strftime('%Y%m%d')
        default_name = f"总组立课重点检查线体_{date_str}.pdf"
        path, _ = QFileDialog.getSaveFileName(self, "保存PDF报告", default_name, "PDF Files (*.pdf)")
        if not path:
            return
        try:
            self.status_bar.showMessage("正在生成PDF...")
            output_path = export_to_pdf(self.current_results, self.selected_file, path)
            self.current_pdf_path = output_path
            self.status_bar.showMessage(f"✅ PDF已保存: {output_path}")
            try:
                if sys.platform == 'darwin':
                    subprocess.run(['open', output_path])
                elif sys.platform == 'win32':
                    os.startfile(output_path)
                else:
                    subprocess.run(['xdg-open', output_path])
            except Exception:
                pass
            QMessageBox.information(self, "成功", f"PDF报告已保存:\n{output_path}")
        except Exception as e:
            QMessageBox.warning(self, "错误", f"导出PDF失败:\n{str(e)}")
            self.status_bar.showMessage("PDF导出失败")
    
    def config_email(self):
        dialog = EmailConfigDialog(self)
        if dialog.exec_():
            dialog.save_config()
            QMessageBox.information(self, "成功", "邮件配置已保存")
    
    def send_email(self):
        if not self.current_results:
            return
        config = load_email_config()
        if not config.get('password'):
            QMessageBox.warning(self, "请先配置", "请先设置邮件配置")
            self.config_email()
            return
        if not config.get('recipient'):
            QMessageBox.warning(self, "请先配置", "请先设置收件人邮箱")
            self.config_email()
            return
        try:
            self.status_bar.showMessage("正在生成PDF...")
            date_str = datetime.now().strftime('%Y%m%d')
            temp_pdf = f"/tmp/总组立课重点检查线体_{date_str}.pdf"
            output_path = export_to_pdf(self.current_results, self.selected_file, temp_pdf)
            self.current_pdf_path = output_path
        except Exception as e:
            QMessageBox.warning(self, "错误", f"生成PDF失败:\n{str(e)}")
            return
        try:
            self.status_bar.showMessage("正在发送邮件...")
            send_email_with_pdf(self.current_pdf_path, config['recipient'])
            self.status_bar.showMessage(f"✅ 邮件已发送至: {config['recipient']}")
            QMessageBox.information(self, "成功", f"邮件已发送至:\n{config['recipient']}")
        except Exception as e:
            QMessageBox.warning(self, "错误", f"发送失败:\n{str(e)}")
            self.status_bar.showMessage("邮件发送失败")


def main():
    try:
        print("Creating QApplication...")
        app = QApplication(sys.argv)
        print("QApplication created")
        app.setStyle('Fusion')
        print("Style set")
        font = QFont()
        font.setFamily('Microsoft YaHei UI')
        app.setFont(font)
        print("Font set")
        print("Creating SelfCheckWindow...")
        window = SelfCheckWindow()
        print("SelfCheckWindow created")
        print("Calling window.show()...")
        window.show()
        print("window.show() completed")
        print("Entering event loop with app.exec_()...")
        result = app.exec_()
        print(f"Event loop exited with code: {result}")
        sys.exit(result)
    except Exception as e:
        print(f"Error in main: {e}")
        import traceback
        traceback.print_exc()
        input("Press Enter to exit...")


if __name__ == '__main__':
    try:
        main()
    except Exception as e:
        print(f"Error in __main__: {e}")
        import traceback
        traceback.print_exc()
        input("Press Enter to exit...")
